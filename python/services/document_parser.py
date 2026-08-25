"""
统一文档解析引擎 — marker-pdf / docling / pymupdf 三级降级

这是文档解析层的现代化核心:
  1. PDF     → marker-pdf (结构化 Markdown, 内置 OCR + 表格) → pymupdf (表格/排版) → pypdf
  2. Office  → docling (docx/pptx 结构化) → python-docx / python-pptx
  3. 图片    → marker OCR (surya) → pytesseract + LLM 视觉理解
  4. 表格    → openpyxl / csv
  5. Markdown→ 标题层级感知解析

设计原则:
  - 所有重型依赖 (marker-pdf / docling) 均为懒加载, 未安装时自动降级
  - 解析结果统一为 ParseBlock 列表, 携带结构化元数据 (heading_path / is_table / is_image / page)
"""

from __future__ import annotations

import os
import re
from dataclasses import dataclass, field
from typing import Any

from schema import DocType

# docx 样式名级别提取 (如 "heading 1" / "标题 2")
_STYLE_LEVEL_RE = re.compile(r"(\d+)")

@dataclass
class ParseBlock:
    """解析后的文档块 — 携带内容 + 结构化元数据"""
    kind: str  # "paragraph" | "heading" | "table" | "image" | "text"
    content: str
    page: int = 0
    heading_path: list[str] = field(default_factory=list)  # 如 ["3.2 架构设计", "3.2.1 为什么选 LangGraph"]
    is_table: bool = False
    is_image: bool = False
    ocr_text: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def heading_str(self) -> str:
        return " / ".join(self.heading_path)


class UnifiedDocumentParser:
    """统一解析入口 — 根据文件类型路由到对应解析器"""

    def __init__(self) -> None:
        self._marker = None
        self._docling = None

    # ── public API ───────────────────────────────────────────

    def parse(self, file_path: str, doc_type: DocType) -> list[ParseBlock]:
        """解析单个文件, 返回结构化块列表"""
        if doc_type == DocType.PDF:
            return self.parse_pdf(file_path)
        if doc_type == DocType.IMAGE:
            return self.parse_image(file_path)
        if doc_type == DocType.TABLE:
            return self.parse_table(file_path)
        if doc_type == DocType.WORD:
            return self.parse_docx(file_path)
        if doc_type == DocType.PPT:
            return self.parse_pptx(file_path)
        if doc_type == DocType.MARKDOWN:
            return self.parse_markdown(file_path)
        return self.parse_text(file_path)

    # ── PDF: marker-pdf → pymupdf → pypdf ────────────────────

    def parse_pdf(self, file_path: str) -> list[ParseBlock]:
        from config import settings
        if settings.marker_enabled:
            blocks = self._parse_pdf_with_marker(file_path)
            if blocks:
                return blocks
        blocks = self._parse_pdf_with_pymupdf(file_path)
        if blocks:
            return blocks
        return self._parse_pdf_with_pypdf(file_path)

    def _parse_pdf_with_marker(self, file_path: str) -> list[ParseBlock]:
        """marker-pdf: PDF → 结构化 Markdown (标题/表格/OCR 全保留)"""
        try:
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            from marker.output import text_from_rendered

            converter = PdfConverter(artifact_dict=create_model_dict())
            rendered = converter(file_path)
            markdown_text, _, _ = text_from_rendered(rendered)
            if not markdown_text or not markdown_text.strip():
                return []
            return self._markdown_to_blocks(markdown_text, page_prefix=1)
        except Exception:
            return []

    def _parse_pdf_with_pymupdf(self, file_path: str) -> list[ParseBlock]:
        """pymupdf 保底: 高质量文本 + 表格提取 + 字号启发式标题识别"""
        try:
            import fitz  # pymupdf

            doc = fitz.open(file_path)
            blocks: list[ParseBlock] = []
            current_heading: list[str] = []

            for page_no in range(len(doc)):
                page = doc[page_no]

                # 1. 表格提取 (pymupdf 原生, 2024 年起稳定可用)
                try:
                    for table in page.find_tables().tables:
                        rows = table.extract() or []
                        if not rows:
                            continue
                        md_rows = []
                        for row in rows:
                            cells = [str(c).strip() if c is not None else "" for c in row]
                            md_rows.append("| " + " | ".join(cells) + " |")
                        if md_rows:
                            header = md_rows[0]
                            md_rows.insert(1, "|" + "---|" * (len(rows[0]) if rows[0] else 1))
                            blocks.append(ParseBlock(
                                kind="table",
                                content=header + "\n" + "\n".join(md_rows[1:]),
                                page=page_no + 1,
                                heading_path=list(current_heading),
                                is_table=True,
                            ))
                except Exception:
                    pass  # 部分页面无表格, 忽略

                # 2. 文本块提取 + 字号启发式标题识别
                text_dict = page.get_text("dict")
                for block in text_dict.get("blocks", []):
                    if block.get("type") != 0:  # 0 = 文本块, 1 = 图片块
                        continue
                    for line in block.get("lines", []):
                        spans = line.get("spans", [])
                        if not spans:
                            continue
                        text = "".join(s.get("text", "") for s in spans).strip()
                        if not text:
                            continue
                        max_size = max(s.get("size", 0) for s in spans)
                        if max_size >= 14.5 and len(text) < 80:
                            # 大字号短文本 → 标题, 由字号推断层级
                            current_heading = self._push_heading(
                                current_heading, text, page_no + 1, font_size=max_size
                            )
                            blocks.append(ParseBlock(
                                kind="heading",
                                content=text,
                                page=page_no + 1,
                                heading_path=list(current_heading),
                            ))
                        else:
                            blocks.append(ParseBlock(
                                kind="paragraph",
                                content=text,
                                page=page_no + 1,
                                heading_path=list(current_heading),
                            ))

            doc.close()
            return blocks
        except Exception:
            return []

    def _parse_pdf_with_pypdf(self, file_path: str) -> list[ParseBlock]:
        """pypdf 最终降级: 纯文本提取"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            blocks: list[ParseBlock] = []
            for page_no, page in enumerate(reader.pages):
                text = (page.extract_text() or "").strip()
                if text:
                    blocks.append(ParseBlock(kind="text", content=text, page=page_no + 1))
            return blocks
        except Exception:
            return [ParseBlock(kind="text", content=f"[PDF 解析失败] {file_path}")]

    # ── Office: docling → python-docx / python-pptx ──────────

    def parse_docx(self, file_path: str) -> list[ParseBlock]:
        blocks = self._parse_with_docling(file_path)
        if blocks:
            return blocks
        return self._parse_docx_with_python_docx(file_path)

    def parse_pptx(self, file_path: str) -> list[ParseBlock]:
        blocks = self._parse_with_docling(file_path)
        if blocks:
            return blocks
        return self._parse_pptx_with_python_pptx(file_path)

    def _parse_with_docling(self, file_path: str) -> list[ParseBlock]:
        """docling: docx/pptx → 结构化 Markdown (IBM 开源)"""
        try:
            from docling.document_converter import DocumentConverter

            converter = DocumentConverter()
            result = converter.convert(file_path)
            markdown_text = result.document.export_to_markdown()
            if not markdown_text or not markdown_text.strip():
                return []
            return self._markdown_to_blocks(markdown_text)
        except Exception:
            return []

    def _parse_docx_with_python_docx(self, file_path: str) -> list[ParseBlock]:
        try:
            import docx

            document = docx.Document(file_path)
            blocks: list[ParseBlock] = []
            current_heading: list[str] = []
            for para in document.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                style = (para.style.name or "").lower()
                if "heading" in style or "标题" in style:
                    # 从样式名提取级别 (如 "Heading 1" → 1), 映射到字号区间复用同一套层级逻辑
                    m = _STYLE_LEVEL_RE.search(style)
                    style_level = int(m.group(1)) if m else 2
                    font_size_proxy = {1: 20.0, 2: 16.0, 3: 13.0}.get(style_level, 16.0)
                    current_heading = self._push_heading(
                        current_heading, text, 0, font_size=font_size_proxy
                    )
                    blocks.append(ParseBlock(kind="heading", content=text, heading_path=list(current_heading)))
                else:
                    blocks.append(ParseBlock(kind="paragraph", content=text, heading_path=list(current_heading)))

            for table in document.tables:
                rows = []
                for row in table.rows:
                    rows.append("| " + " | ".join(cell.text.strip() for cell in row.cells) + " |")
                if rows:
                    blocks.append(ParseBlock(kind="table", content="\n".join(rows), is_table=True))
            return blocks or [ParseBlock(kind="text", content=f"[DOCX 解析失败] {file_path}")]
        except Exception:
            return [ParseBlock(kind="text", content=f"[DOCX 解析失败] {file_path}")]

    def _parse_pptx_with_python_pptx(self, file_path: str) -> list[ParseBlock]:
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            blocks: list[ParseBlock] = []
            for slide_no, slide in enumerate(prs.slides, start=1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(run.text for run in para.runs).strip()
                            if t:
                                texts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            texts.append(" | ".join(cell.text.strip() for cell in row.cells))
                if texts:
                    blocks.append(ParseBlock(
                        kind="paragraph",
                        content=f"[Slide {slide_no}]\n" + "\n".join(texts),
                        page=slide_no,
                    ))
            return blocks or [ParseBlock(kind="text", content=f"[PPTX 解析失败] {file_path}")]
        except Exception:
            return [ParseBlock(kind="text", content=f"[PPTX 解析失败] {file_path}")]

    # ── 图片: marker OCR → pytesseract ───────────────────────

    def parse_image(self, file_path: str) -> list[ParseBlock]:
        from config import settings
        blocks = self._parse_image_with_marker(file_path) if settings.marker_enabled else []
        if blocks:
            return blocks

        # 降级: pytesseract OCR (无 tesseract 时返回空, 由调用方走 LLM 视觉)
        ocr_text = self._ocr_with_tesseract(file_path)
        if ocr_text.strip():
            return [ParseBlock(kind="text", content=ocr_text, is_image=True, ocr_text=ocr_text)]
        return [ParseBlock(kind="image", content="", is_image=True, metadata={"needs_vision": True})]

    def _parse_image_with_marker(self, file_path: str) -> list[ParseBlock]:
        try:
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            from marker.output import text_from_rendered

            converter = PdfConverter(artifact_dict=create_model_dict())
            rendered = converter(file_path)
            markdown_text, _, _ = text_from_rendered(rendered)
            if not markdown_text or not markdown_text.strip():
                return []
            return self._markdown_to_blocks(markdown_text)
        except Exception:
            return []

    @staticmethod
    def _ocr_with_tesseract(file_path: str) -> str:
        try:
            import pytesseract
            from PIL import Image
            return pytesseract.image_to_string(Image.open(file_path), lang="chi_sim+eng")
        except Exception:
            return ""

    # ── 表格: openpyxl / csv ─────────────────────────────────

    @staticmethod
    def parse_table(file_path: str) -> list[ParseBlock]:
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".csv":
                return UnifiedDocumentParser._parse_csv(file_path)
            return UnifiedDocumentParser._parse_excel(file_path)
        except Exception:
            return [ParseBlock(kind="text", content=f"[表格解析失败] {file_path}")]

    @staticmethod
    def _parse_csv(file_path: str) -> list[ParseBlock]:
        import csv
        blocks: list[ParseBlock] = []
        with open(file_path, encoding="utf-8") as f:
            reader = csv.DictReader(f)
            headers = reader.fieldnames or []
            rows: list[str] = []
            for row in reader:
                rows.append("| " + " | ".join(str(row.get(h, "")) for h in headers) + " |")
            if rows:
                blocks.append(ParseBlock(
                    kind="table",
                    content="| " + " | ".join(headers) + " |\n|---" * len(headers) + "\n" + "\n".join(rows),
                    is_table=True,
                ))
        return blocks or [ParseBlock(kind="text", content="[空 CSV]")]

    @staticmethod
    def _parse_excel(file_path: str) -> list[ParseBlock]:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True)
        blocks: list[ParseBlock] = []
        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue
            headers = [str(c) if c is not None else "" for c in rows[0]]
            md_rows = ["| " + " | ".join(headers) + " |", "|---" * len(headers)]
            for row in rows[1:]:
                md_rows.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
            blocks.append(ParseBlock(
                kind="table",
                content=f"[工作表: {sheet.title}]\n" + "\n".join(md_rows),
                is_table=True,
            ))
        return blocks or [ParseBlock(kind="text", content="[空 Excel]")]

    # ── Markdown / 纯文本 ────────────────────────────────────

    @staticmethod
    def parse_markdown(file_path: str) -> list[ParseBlock]:
        with open(file_path, encoding="utf-8") as f:
            return UnifiedDocumentParser._markdown_to_blocks(f.read())

    @staticmethod
    def parse_text(file_path: str) -> list[ParseBlock]:
        with open(file_path, encoding="utf-8") as f:
            content = f.read()
        return [ParseBlock(kind="text", content=content)]

    @staticmethod
    def _markdown_to_blocks(markdown_text: str, page_prefix: int = 0) -> list[ParseBlock]:
        """Markdown → 结构化块, 跟踪标题层级 (heading_path)"""
        blocks: list[ParseBlock] = []
        current_path: list[str] = []
        current_content: list[str] = []
        current_kind = "paragraph"
        current_is_table = False

        def flush() -> None:
            nonlocal current_content, current_kind, current_is_table
            if current_content:
                text = "\n".join(current_content).strip()
                if text:
                    blocks.append(ParseBlock(
                        kind=current_kind,
                        content=text,
                        page=page_prefix,
                        heading_path=list(current_path),
                        is_table=current_is_table,
                    ))
            current_content = []
            current_is_table = False

        for line in markdown_text.splitlines():
            stripped = line.strip()
            if not stripped:
                continue

            # 表格行
            if stripped.startswith("|"):
                if not current_is_table:
                    flush()
                    current_kind = "table"
                    current_is_table = True
                current_content.append(stripped)
                continue
            if current_is_table:
                flush()
                current_kind = "paragraph"

            # 标题行
            if stripped.startswith("#"):
                flush()
                level = len(stripped) - len(stripped.lstrip("#"))
                title = stripped.lstrip("#").strip()
                current_path = current_path[:level - 1] + [title]
                blocks.append(ParseBlock(
                    kind="heading",
                    content=title,
                    page=page_prefix,
                    heading_path=list(current_path),
                ))
                current_kind = "paragraph"
                continue

            current_content.append(stripped)

        flush()
        return blocks

    @staticmethod
    def _push_heading(current: list[str], new_heading: str, page: int, font_size: float = 0) -> list[str]:
        """
        根据字号推断标题层级, 替换同层及以下标题:
          size >= 18   → level 1
          size >= 14.5 → level 2
          size >= 12   → level 3
          其他         → 非标题 (不修改 current)
        """
        if font_size >= 18:
            level = 1
        elif font_size >= 14.5:
            level = 2
        elif font_size >= 12:
            level = 3
        else:
            # 不够大，视为普通文本，不更新 heading_path
            return list(current)

        if level == 1:
            return [new_heading]
        # 保留上一层级之前的所有标题，替换当前层级
        return current[:level - 1] + [new_heading]
